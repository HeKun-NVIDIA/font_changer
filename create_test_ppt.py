#!/usr/bin/env python3
"""
创建测试PPT文件，包含中英文混合文本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_test_ppt():
    """创建包含中英文混合文本的测试PPT"""
    
    # 创建新的演示文稿
    prs = Presentation()
    
    # 第一张幻灯片 - 标题页
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # 标题布局
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "PPT字体替换测试 Font Replacement Test"
    subtitle.text = "这是一个测试文件 This is a test file\n包含中英文混合内容 Contains mixed Chinese and English content"
    
    # 第二张幻灯片 - 内容页
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
    title2 = slide2.shapes.title
    content = slide2.placeholders[1]
    
    title2.text = "测试内容 Test Content"
    
    # 添加混合文本内容
    tf = content.text_frame
    tf.text = "这是第一段中文内容。This is the first English paragraph."
    
    p = tf.add_paragraph()
    p.text = "Hello 世界！这里有数字123和符号@#$%。"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "AI人工智能(Artificial Intelligence)正在改变世界。"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "iPhone 15 Pro Max售价￥9999，性能提升50%。"
    p.level = 1
    
    # 第三张幻灯片 - 添加文本框
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 添加标题文本框
    title_box = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "自定义文本框测试 Custom TextBox Test"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # 添加内容文本框
    content_box = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    
    # 添加多个段落
    p1 = content_frame.paragraphs[0]
    p1.text = "这是一个包含中英文的段落。This paragraph contains both Chinese and English."
    
    p2 = content_frame.add_paragraph()
    p2.text = "Python编程语言 + JavaScript = 全栈开发"
    
    p3 = content_frame.add_paragraph()
    p3.text = "2024年，ChatGPT用户突破1亿。In 2024, ChatGPT users exceeded 100 million."
    
    # 第四张幻灯片 - 表格测试
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # 标题和内容布局
    title4 = slide4.shapes.title
    title4.text = "表格测试 Table Test"
    
    # 添加表格
    rows, cols = 3, 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3)
    
    table = slide4.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 填充表格内容
    table.cell(0, 0).text = "项目 Item"
    table.cell(0, 1).text = "中文名称"
    table.cell(0, 2).text = "English Name"
    
    table.cell(1, 0).text = "产品A"
    table.cell(1, 1).text = "智能手机"
    table.cell(1, 2).text = "Smartphone"
    
    table.cell(2, 0).text = "产品B"
    table.cell(2, 1).text = "笔记本电脑"
    table.cell(2, 2).text = "Laptop Computer"
    
    # 保存文件
    output_path = "/home/ubuntu/ppt-font-changer/test_presentation.pptx"
    prs.save(output_path)
    print(f"测试PPT文件已创建: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_test_ppt()

